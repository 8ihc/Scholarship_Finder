"""
Document Parsing and OCR Staging Script
This script processes downloaded scholarship application documents,
attempting native text extraction first and flagging documents for OCR if necessary.
"""

# document_parser.py (Optimized for ID_Name.ext filename format and includes image types)
import os
import re
import json
import pandas as pd 
from pathlib import Path
from typing import Tuple
from docx import Document
from odf.opendocument import load
from odf.text import P, Span
from odf import teletype


# Document Libraries
# NOTE: Ensure these are installed via requirements.txt
try:
    import pdfplumber
    from docx import Document
    from odf.opendocument import load
    from odf import text
    import openpyxl # For XLSX/XLS/ODS
    from pptx import Presentation # For PPTX
except ImportError:
    print("WARNING: Document libraries not fully installed. Check requirements.txt.")

# --- Configuration ---
SCHOLARSHIPS_FILE = Path("data/raw/scholarships.csv") 
ATTACHMENTS_DIR = Path("data/raw/attachments")
PARSED_OUTPUT_FILE = Path("data/processed/scholarships_parsed_texts.json")
# ---------------------

# Regex to extract ID and Name from the filename (e.g., 7652_專用申請書.pdf)
FILENAME_PATTERN = re.compile(r'(\d+)_(.+?)\.([a-zA-Z0-9]+)$')

# 常見圖像副檔名
IMAGE_EXTENSIONS = ['png', 'jpg', 'jpeg']


# --- Core Parsing Functions ---

def extract_text_from_pdf(pdf_path: Path) -> Tuple[str, bool]:
    """
    Extracts text from PDF (Pass 1: Native Extraction). 
    Returns text and a boolean indicating if OCR is REQUIRED.
    """
    all_text = []
    ocr_needed = False
    
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text(x_tolerance=2)
                
                if text and text.strip():
                    all_text.append(text)
                else:
                    # Native extraction failed, marking for OCR
                    all_text.append(f"[PAGE_OCR_REQUIRED]")
                    ocr_needed = True
                    
        return "\n".join(all_text), ocr_needed

    except Exception as e:
        return f"[ERROR_PARSING_PDF: {e}]", False

def extract_text_from_word_odt(file_path: Path) -> str:
    """Extracts text from DOCX/DOC/ODT files."""
    ext = file_path.suffix.lower()
    try:
        if ext == '.docx':
            doc = Document(file_path)
            return '\n'.join([p.text for p in doc.paragraphs])

        elif ext == '.odt':
            # 正確使用 odfpy 讀取
            doc = load(file_path)

            # 使用 odfpy 的 teletype 模組安全地提取純文字
            # teletype.extractText 會自動遍歷 ODF 的元素
            text_content = teletype.extractText(doc.text)
            return text_content.strip()
            
        elif ext == '.doc':
            # .doc files (old binary format) require external conversion (e.g., LibreOffice)
            return "[DOC_OLD_FORMAT: Requires conversion or external tools]"
        
        return f"[UNSUPPORTED_FORMAT: {ext}]"
        
    except Exception as e:
        return f"[ERROR_PARSING_{ext.upper()}: {e}]"

def extract_text_from_excel(file_path: Path) -> str:
    """Extracts text from XLSX/XLS/ODS files."""
    ext = file_path.suffix.lower()
    all_text = []
    try:
        if ext == '.xlsx':
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in workbook.worksheets:
                all_text.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows():
                    row_data = [str(cell.value) if cell.value is not None else '' for cell in row]
                    all_text.append("|".join(row_data))
            return "\n".join(all_text)
        # Note: ODS parsing needs the ODF logic above, or specialized ODS Excel parsing
        return f"[UNSUPPORTED_EXCEL_FORMAT: {ext}]"
    except Exception as e:
        return f"[ERROR_PARSING_EXCEL: {e}]"

def extract_text_from_ppt(file_path: Path) -> str:
    """Extracts text from PPTX files."""
    all_text = []
    try:
        presentation = Presentation(file_path)
        for slide_index, slide in enumerate(presentation.slides):
            all_text.append(f"\n--- Slide {slide_index + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text)
    except Exception as e:
        return f"[ERROR_PARSING_PPTX: {e}]"


def process_all_files():
    """
    Main loop to process all downloaded files.
    """
    
    if not ATTACHMENTS_DIR.exists():
        print(f"Error: Attachment directory not found at {ATTACHMENTS_DIR}")
        return
        
    parsed_results = []
    
    print(f"Scanning directory: {ATTACHMENTS_DIR}")
    
    for file_path in ATTACHMENTS_DIR.iterdir():
        if not file_path.is_file():
            continue
        
        # --- 核心新增邏輯: 忽略備份檔案 ---
        if file_path.name.lower().endswith('_archive'):
            print(f"Skipping archive file: {file_path.name}")
            continue
        # --- 核心新增邏輯結束 ---
            
        # 1. Extract Metadata from Filename
        match = FILENAME_PATTERN.search(file_path.name)
        if not match:
            print(f"Skipping file due to unexpected naming format: {file_path.name}")
            continue

        sch_id = match.group(1) 
        name = match.group(2).strip() 
        ext = match.group(3).lower() 
        
        print(f"\nProcessing ID {sch_id} - Attachment: {name}.{ext}")

        # 2. Execute Two-Pass Parsing / Image Flagging
        ocr_required = False
        parsed_text = ""
        status = 'ERROR_UNKNOWN' # Default error status

        # --- 處理圖像檔案 (.png, .jpg) ---
        if ext in IMAGE_EXTENSIONS:
            parsed_text = f"[IMAGE_FILE: .'{ext.upper()}' needs Cloud Vision OCR]"
            ocr_required = True
        
        # --- 處理 PDF ---
        elif ext == 'pdf':
            parsed_text, ocr_required = extract_text_from_pdf(file_path)
        
        # --- 處理 Word/ODT ---
        elif ext in ['docx', 'doc', 'odt']:
            parsed_text = extract_text_from_word_odt(file_path)
        
        # --- 處理試算表 (.xlsx, .ods) ---
        elif ext in ['xlsx', 'ods']:
            parsed_text = extract_text_from_excel(file_path)
        
        # --- 處理簡報 (.pptx) ---
        elif ext == 'pptx':
            parsed_text = extract_text_from_ppt(file_path)

        # --- 處理其他未支援格式 ---
        else:
            parsed_text = f"[UNSUPPORTED_TYPE: .{ext}]"

        # 3. Determine Final Status
        status = 'OCR_REQUIRED' if ocr_required else 'NATIVE_OK'
        
        # 覆寫狀態：如果解析失敗 (Placeholder 或 ERROR)
        if 'ERROR' in parsed_text or 'DOC_OLD_FORMAT' in parsed_text or 'UNSUPPORTED_TYPE' in parsed_text: 
            status = 'ERROR_PARSE'

        # 4. Record Result
        parsed_results.append({
            'id': sch_id,
            'name': name,
            'status': status,
            'file_path_local': str(file_path),
            'parsed_text': parsed_text,
        })
        
        print(f"  -> Status: {status}")

    # 5. Save results to JSON
    PARSED_OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(PARSED_OUTPUT_FILE, 'w', encoding='utf-8') as f:
        json.dump(parsed_results, f, ensure_ascii=False, indent=4)
        
    print(f"\n--- PARSING COMPLETE ---")
    print(f"Results saved to: {PARSED_OUTPUT_FILE}")
    
    # 6. Report Summary
    ocr_count = sum(1 for r in parsed_results if r['status'] == 'OCR_REQUIRED')
    native_count = sum(1 for r in parsed_results if r['status'] == 'NATIVE_OK')
    error_count = sum(1 for r in parsed_results if r['status'] == 'ERROR_PARSE')
    
    print(f"Total files processed: {len(parsed_results)}")
    print(f"Native Parsing Success: {native_count} files")
    print(f"❌ OCR Required:        {ocr_count} files (包括圖像檔案和掃描 PDF)")
    print(f"⚠️ Errors/Unsupported:   {error_count} files")
    
    return parsed_results


if __name__ == '__main__':
    process_all_files()